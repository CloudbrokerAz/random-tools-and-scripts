#!/usr/bin/env python3
"""
build_presentation.py -- Programmatically build IBM-branded PowerPoint presentations.

Uses the IBM POTX template (IBM_presentation_brand_covers_v_2_1_Plex_embed.potx)
and a JSON specification to generate PPTX files with full support for all 49
slide layouts.

Usage:
    python build_presentation.py spec.json
    python build_presentation.py spec.json -o output.pptx
    cat spec.json | python build_presentation.py -
    python build_presentation.py --list-layouts
"""

from __future__ import annotations

import argparse
import copy
import json
import os
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# IBM Carbon Design System Colors
# ---------------------------------------------------------------------------
IBM_BLUE_60 = RGBColor(0x0F, 0x62, 0xFE)
IBM_PURPLE_50 = RGBColor(0xA5, 0x6E, 0xFF)
IBM_CYAN_80 = RGBColor(0x00, 0x3A, 0x6D)
IBM_TEAL_50 = RGBColor(0x00, 0x9D, 0x9A)
IBM_MAGENTA_70 = RGBColor(0x9F, 0x18, 0x53)
IBM_RED_50 = RGBColor(0xFA, 0x4D, 0x56)
IBM_CYAN_10 = RGBColor(0xE5, 0xF6, 0xFF)
IBM_CYAN_20 = RGBColor(0xBA, 0xE6, 0xFF)
IBM_GRAY_10 = RGBColor(0xF4, 0xF4, 0xF4)
IBM_GRAY_100 = RGBColor(0x16, 0x16, 0x16)
IBM_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# HashiCorp product colors (IBM palette equivalents)
TERRAFORM_COLOR = IBM_PURPLE_50
VAULT_COLOR = RGBColor(0xF1, 0xC2, 0x1B)  # Yellow 30
CONSUL_COLOR = RGBColor(0xEE, 0x53, 0x96)  # Magenta 50
NOMAD_COLOR = RGBColor(0x24, 0xA1, 0x48)   # Green 50
PACKER_COLOR = RGBColor(0x11, 0x92, 0xE8)  # Cyan 50
BOUNDARY_COLOR = IBM_RED_50
WAYPOINT_COLOR = IBM_TEAL_50
VAGRANT_COLOR = IBM_BLUE_60

# Default font
IBM_FONT = "IBM Plex Sans"
IBM_FONT_LIGHT = "IBM Plex Sans Light"

# ---------------------------------------------------------------------------
# Layout metadata: maps layout index -> name and ordered placeholder semantics
# ---------------------------------------------------------------------------
# Each entry: { "name": str, "placeholders": { idx: semantic_role } }
# Semantic roles used:
#   title, subtitle, body, body2, body_right, image, image_right,
#   footer, slide_number, data_value_N, data_label_N, profile_image_N,
#   profile_info_N, icon_N, table, chart_title, chart_body, legal_body,
#   legal_body2

LAYOUT_META: dict[int, dict[str, Any]] = {
    0: {
        "name": "Cover, imagery",
        "category": "cover",
        "placeholders": {0: "title", 11: "image", 12: "body", 13: "subtitle"},
    },
    1: {
        "name": "Cover, cyan",
        "category": "cover",
        "placeholders": {0: "title", 12: "body", 13: "subtitle"},
    },
    2: {
        "name": "Cover, plain",
        "category": "cover",
        "placeholders": {0: "title"},
    },
    3: {
        "name": "Cover, plain, label",
        "category": "cover",
        "placeholders": {0: "title", 10: "subtitle"},
    },
    4: {
        "name": "Cover, imagery, half",
        "category": "cover",
        "placeholders": {0: "title", 12: "image"},
    },
    5: {
        "name": "Cover, imagery, half, label",
        "category": "cover",
        "placeholders": {0: "title", 11: "subtitle", 12: "image"},
    },
    6: {
        "name": "Contents",
        "category": "navigation",
        "placeholders": {
            0: "title", 10: "body", 11: "body_right",
            18: "footer", 4: "slide_number",
        },
    },
    7: {
        "name": "Section divider",
        "category": "navigation",
        "placeholders": {0: "title", 18: "footer", 4: "slide_number"},
    },
    8: {
        "name": "Large text",
        "category": "navigation",
        "placeholders": {0: "title", 18: "footer", 4: "slide_number"},
    },
    9: {
        "name": "Callout, headline",
        "category": "callout",
        "placeholders": {
            0: "title", 10: "body", 18: "footer", 4: "slide_number",
        },
    },
    10: {
        "name": "Callout, stand-alone",
        "category": "callout",
        "placeholders": {19: "body", 18: "footer", 4: "slide_number"},
    },
    11: {
        "name": "Data, 2 callouts, vertical",
        "category": "data",
        "placeholders": {
            11: "data_value_1", 12: "data_label_1",
            13: "data_value_2", 14: "data_label_2",
            18: "footer", 4: "slide_number",
        },
    },
    12: {
        "name": "Data, 3 callouts, vertical",
        "category": "data",
        "placeholders": {
            0: "title",
            15: "data_label_1", 12: "data_value_1",
            16: "data_label_2", 13: "data_value_2",
            17: "data_label_3", 14: "data_value_3",
            11: "body",
            18: "footer", 4: "slide_number",
        },
    },
    13: {
        "name": "Data, 2 callouts, horizontal",
        "category": "data",
        "placeholders": {
            0: "title",
            11: "data_value_1", 12: "data_label_1",
            13: "data_value_2",
            18: "footer", 4: "slide_number",
        },
    },
    14: {
        "name": "Data, 3 callouts, horizontal",
        "category": "data",
        "placeholders": {
            0: "title",
            15: "data_label_1", 12: "data_value_1",
            16: "data_label_2", 13: "data_value_2",
            17: "data_label_3", 14: "data_value_3",
            11: "body",
            18: "footer", 4: "slide_number",
        },
    },
    15: {
        "name": "Text, 4 columns",
        "category": "text_columns",
        "placeholders": {
            0: "title",
            11: "col_1", 12: "col_2", 13: "col_3", 14: "col_4",
            18: "footer", 4: "slide_number",
        },
    },
    16: {
        "name": "Text, 4 columns, short dividers",
        "category": "text_columns",
        "placeholders": {
            0: "title", 15: "subtitle",
            11: "col_1", 12: "col_2", 13: "col_3", 14: "col_4",
            18: "footer", 4: "slide_number",
        },
    },
    17: {
        "name": "Text, 4 columns, dividers, headlines",
        "category": "text_columns",
        "placeholders": {
            0: "title", 15: "subtitle",
            12: "col_1", 16: "col_1_head",
            13: "col_2", 17: "col_2_head",
            14: "col_3",
            18: "footer", 4: "slide_number",
        },
    },
    18: {
        "name": "Text, 4 columns, dividers, pictograms",
        "category": "text_columns",
        "placeholders": {
            0: "title",
            11: "col_1", 15: "icon_1",
            12: "col_2", 16: "icon_2",
            13: "col_3", 17: "icon_3",
            14: "col_4",
            18: "footer", 4: "slide_number",
        },
    },
    19: {
        "name": "Text, 1 wide column, divider",
        "category": "text_columns",
        "placeholders": {
            0: "title", 15: "subtitle",
            12: "body", 13: "body_right",
            18: "footer", 4: "slide_number",
        },
    },
    20: {
        "name": "Text, 2 wide columns",
        "category": "text_columns",
        "placeholders": {
            0: "title", 11: "body", 12: "body_right",
            18: "footer", 4: "slide_number",
        },
    },
    21: {
        "name": "Text, 2 columns, large title",
        "category": "text_columns",
        "placeholders": {
            0: "title", 11: "body", 12: "body_right",
            18: "footer", 4: "slide_number",
        },
    },
    22: {
        "name": "Text, 2 columns, small title",
        "category": "text_columns",
        "placeholders": {
            0: "title", 11: "body", 12: "body_right",
            18: "footer", 4: "slide_number",
        },
    },
    23: {
        "name": "Text, 2 columns, dividers, large title",
        "category": "text_columns",
        "placeholders": {
            0: "title",
            16: "col_1_head", 18: "col_1",
            17: "col_2_head", 19: "col_2",
            20: "footer", 4: "slide_number",
        },
    },
    24: {
        "name": "Text, 2 columns, dividers, small title",
        "category": "text_columns",
        "placeholders": {
            0: "title",
            16: "col_1_head", 18: "col_1",
            17: "col_2_head", 19: "col_2",
            20: "footer", 4: "slide_number",
        },
    },
    25: {
        "name": "Text, 2 columns, dividers, pictograms",
        "category": "text_columns",
        "placeholders": {
            0: "title", 15: "subtitle",
            12: "col_1", 16: "icon_1",
            13: "col_2", 17: "icon_2",
            14: "body",
            18: "footer", 4: "slide_number",
        },
    },
    26: {
        "name": "Boxes, 4 stacked wide, pictograms",
        "category": "box_grid",
        "placeholders": {
            11: "box_1", 15: "icon_1",
            12: "box_2", 16: "icon_2",
            13: "box_3", 17: "icon_3",
            14: "box_4", 18: "icon_4",
            19: "footer", 4: "slide_number",
        },
    },
    27: {
        "name": "Boxes, 4 stacked, small title",
        "category": "box_grid",
        "placeholders": {
            0: "title", 21: "subtitle",
            16: "box_1_head", 18: "box_1",
            17: "box_2_head", 19: "box_2",
            22: "footer", 4: "slide_number",
        },
    },
    28: {
        "name": "Boxes, 4 stacked, large title",
        "category": "box_grid",
        "placeholders": {
            0: "title", 21: "subtitle",
            16: "box_1_head", 18: "box_1",
            17: "box_2_head", 19: "box_2",
            22: "footer", 4: "slide_number",
        },
    },
    29: {
        "name": "Boxes, 4 horizontal, small title",
        "category": "box_grid",
        "placeholders": {
            21: "subtitle",
            11: "box_1", 12: "box_2", 13: "box_3", 14: "box_4",
            18: "footer", 4: "slide_number",
        },
    },
    30: {
        "name": "Boxes, 4 horizontal, large title",
        "category": "box_grid",
        "placeholders": {
            0: "title",
            11: "box_1", 12: "box_2", 13: "box_3", 14: "box_4",
            18: "footer", 4: "slide_number",
        },
    },
    31: {
        "name": "Boxes, 6 stacked",
        "category": "box_grid",
        "placeholders": {
            0: "title",
            20: "box_1_head", 12: "box_1",
            16: "box_2_head", 13: "box_2",
            17: "box_3_head", 14: "box_3",
            22: "footer", 4: "slide_number",
        },
    },
    32: {
        "name": "Boxes, 6 stacked, icons",
        "category": "box_grid",
        "placeholders": {
            0: "title",
            11: "box_1", 17: "icon_1",
            12: "box_2", 18: "icon_2",
            13: "box_3", 19: "icon_3",
            14: "box_4", 20: "icon_4",
            15: "box_5", 21: "icon_5",
            16: "box_6", 22: "icon_6",
            23: "footer", 4: "slide_number",
        },
    },
    33: {
        "name": "Boxes, 6 stacked, alternate, large title",
        "category": "box_grid",
        "placeholders": {
            0: "title",
            11: "box_1", 12: "box_1_head",
            19: "box_2", 13: "box_2_head",
            20: "box_3", 14: "box_4",
            18: "footer", 4: "slide_number",
        },
    },
    34: {
        "name": "Boxes, 6 stacked, alternate, small title",
        "category": "box_grid",
        "placeholders": {
            0: "title",
            11: "box_1", 12: "box_1_head",
            19: "box_2", 13: "box_2_head",
            20: "box_3", 14: "box_4",
            18: "footer", 4: "slide_number",
        },
    },
    35: {
        "name": "Video or imagery, half, inset",
        "category": "media",
        "placeholders": {
            0: "title", 12: "body", 13: "subtitle",
            15: "image", 18: "footer",
        },
    },
    36: {
        "name": "Video or imagery, 3/4, bleed",
        "category": "media",
        "placeholders": {
            0: "title", 13: "body", 14: "image",
            18: "footer", 4: "slide_number",
        },
    },
    37: {
        "name": "Video or imagery, 3/4, inset",
        "category": "media",
        "placeholders": {
            0: "title", 12: "body", 14: "image", 18: "footer",
        },
    },
    38: {
        "name": "Video or imagery, bleed",
        "category": "media",
        "placeholders": {11: "image", 4: "slide_number"},
    },
    39: {
        "name": "Video or imagery, inset",
        "category": "media",
        "placeholders": {11: "image"},
    },
    40: {
        "name": "Contacts, profiles, contributors",
        "category": "special",
        "placeholders": {
            0: "title",
            # 6 profile slots: image placeholder + text placeholder pairs
            28: "profile_image_1", 21: "profile_info_1",
            30: "profile_image_2", 12: "profile_info_2",
            32: "profile_image_3", 14: "profile_info_3",
            34: "profile_image_4", 27: "profile_info_4",
            36: "profile_image_5", 23: "profile_info_5",
            38: "profile_image_6", 25: "profile_info_6",
            18: "footer", 4: "slide_number",
        },
    },
    41: {
        "name": "Table",
        "category": "special",
        "placeholders": {
            0: "title", 11: "table", 18: "footer", 4: "slide_number",
        },
    },
    42: {
        "name": "Chart",
        "category": "special",
        "placeholders": {
            0: "title", 13: "chart_title", 14: "chart_body",
            18: "footer", 4: "slide_number",
        },
    },
    43: {
        "name": "Legal disclaimer, one column",
        "category": "special",
        "placeholders": {
            14: "legal_title", 15: "legal_body",
            16: "footer", 4: "slide_number",
        },
    },
    44: {
        "name": "Legal disclaimer, two columns",
        "category": "special",
        "placeholders": {
            14: "legal_title", 15: "legal_body", 17: "legal_body2",
            18: "footer", 4: "slide_number",
        },
    },
    45: {
        "name": "Blank slide",
        "category": "special",
        "placeholders": {0: "title", 12: "footer", 4: "slide_number"},
    },
    46: {
        "name": "Blank slide, no footer",
        "category": "special",
        "placeholders": {0: "title", 4: "slide_number"},
    },
    47: {
        "name": "End slide",
        "category": "special",
        "placeholders": {},
    },
    48: {
        "name": "DEFAULT",
        "category": "special",
        "placeholders": {},
    },
}


# ---------------------------------------------------------------------------
# Template loading -- POTX requires content-type patching
# ---------------------------------------------------------------------------

def _resolve_template_path() -> Path:
    """Resolve the IBM POTX template path relative to this script."""
    script_dir = Path(__file__).resolve().parent
    skill_dir = script_dir.parent
    template = skill_dir / "resources" / "templates" / "IBM_presentation_brand_covers_v_2_1_Plex_embed.potx"
    if not template.exists():
        raise FileNotFoundError(f"IBM template not found at: {template}")
    return template


def load_template(template_path: Optional[Path] = None) -> Presentation:
    """Load the IBM POTX template, converting it to PPTX content-type on the fly.

    python-pptx refuses to open .potx files directly because of the content-type
    header in [Content_Types].xml. We copy the file, patch the content-type in
    memory, and open the result.
    """
    if template_path is None:
        template_path = _resolve_template_path()

    tmp_pptx = tempfile.mktemp(suffix=".pptx")
    try:
        with zipfile.ZipFile(str(template_path), "r") as zin:
            ct_xml = zin.read("[Content_Types].xml")
            ct_xml = ct_xml.replace(
                b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
            )
            with zipfile.ZipFile(tmp_pptx, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        data = ct_xml
                    zout.writestr(item, data)

        prs = Presentation(tmp_pptx)
        return prs
    finally:
        try:
            os.unlink(tmp_pptx)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# Helper: parse a hex color string to RGBColor
# ---------------------------------------------------------------------------

def parse_color(color_str: str) -> RGBColor:
    """Parse '#RRGGBB' or 'RRGGBB' to an RGBColor."""
    color_str = color_str.lstrip("#")
    if len(color_str) != 6:
        raise ValueError(f"Invalid color: {color_str}")
    r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
    return RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Visual enhancement helpers
# ---------------------------------------------------------------------------

def _add_accent_bar(slide, x, y, width, height, color):
    """Add a solid colored accent bar (rectangle with no border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _apply_fill_opacity(shape, opacity_pct):
    """Set the fill opacity on a shape via XML manipulation.

    Args:
        shape: A python-pptx shape with a solid fill already applied.
        opacity_pct: Opacity as a percentage (0=fully transparent, 100=fully opaque).
    """
    # Find the solidFill element and add an alpha child
    spPr = shape._element.spPr
    solid_fill = spPr.find(qn('a:solidFill'))
    if solid_fill is not None:
        color_elem = solid_fill.find(qn('a:srgbClr'))
        if color_elem is not None:
            # Remove existing alpha if any
            for existing in color_elem.findall(qn('a:alpha')):
                color_elem.remove(existing)
            alpha = etree.SubElement(color_elem, qn('a:alpha'))
            # Value is in 1000ths of a percent (100% = 100000)
            alpha.set('val', str(int(opacity_pct * 1000)))


def _add_scrim_overlay(slide, x, y, width, height, color="#000000", opacity=50):
    """Add a semi-transparent dark overlay for text contrast over images.

    Args:
        slide: The slide to add the scrim to.
        x, y, width, height: Position and size in inches.
        color: Scrim color hex (default black).
        opacity: Opacity percentage (default 50 = 50% opaque).
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_color(color)
    shape.line.fill.background()
    _apply_fill_opacity(shape, opacity)
    return shape


NAMED_BACKGROUNDS = {
    "white": "#FFFFFF",
    "cyan_10": "#E5F6FF",
    "cyan_20": "#BAE6FF",
    "gray_10": "#F4F4F4",
    "gray_100": "#161616",
    "blue_90": "#001D6C",
}


def _set_slide_background(slide, color_spec):
    """Set slide background to a solid color. Accepts named colors or hex."""
    color_hex = NAMED_BACKGROUNDS.get(color_spec, color_spec)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = parse_color(color_hex)


def _add_overlay_images(slide, overlays, base_dir):
    """Add freeform images at specified absolute positions on the slide."""
    for overlay in overlays:
        try:
            img_path = _resolve_image_path(overlay["image"], base_dir)
            slide.shapes.add_picture(
                str(img_path),
                Inches(overlay.get("x", 0)),
                Inches(overlay.get("y", 0)),
                Inches(overlay.get("width", 1)),
                Inches(overlay.get("height", 1))
            )
        except Exception as e:
            print(f"  WARNING: Could not add overlay image {overlay.get('image')}: {e}", file=sys.stderr)


IBM_GRAY_20 = RGBColor(0xE0, 0xE0, 0xE0)


def _add_divider_line(slide, x, y, length, orientation="vertical", color=None):
    """Add a thin 1pt divider line."""
    if color is None:
        color = IBM_GRAY_20
    elif isinstance(color, str):
        color = parse_color(color)
    if orientation == "vertical":
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Pt(1), Inches(length)
        )
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(length), Pt(1)
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _apply_shadow(shape, blur=50800, dist=25400, direction=5400000, opacity=20):
    """Apply an outer drop shadow to a shape via direct XML manipulation.

    Args:
        shape: A python-pptx shape object
        blur: Shadow blur radius in EMUs (default 50800 = ~4pt)
        dist: Shadow distance in EMUs (default 25400 = ~2pt)
        direction: Shadow angle in 60000ths of a degree (default 5400000 = 90°/bottom)
        opacity: Shadow opacity as percentage (default 20 = 20%)
    """
    spPr = shape._element.spPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', '000000')
    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha.set('val', str(opacity * 1000))  # Convert percentage to 1000ths


def _add_callout_shape(slide, x, y, width, height, fill_color, border_color=None, text=None,
                       rich_text=None, shadow=False, font_size=None, valign=None,
                       text_color=None, corner_radius=None):
    """Add a rounded rectangle callout/highlight shape.

    Args:
        slide: The slide to add the shape to.
        x, y, width, height: Position and size in inches.
        fill_color: Background fill color (hex string or RGBColor).
        border_color: Optional border color (hex string or RGBColor).
        text: Optional plain text content.
        rich_text: Optional list of {"text": str, "bold": bool, "color": str} segments.
        shadow: Whether to apply a drop shadow (default False).
        font_size: Override font size in pt (default 11).
        valign: Vertical alignment - "top", "middle", or "bottom".
        text_color: Hex color for plain text (default None = black).
        corner_radius: Corner radius in inches (default None = shape default).
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color if isinstance(fill_color, RGBColor) else parse_color(fill_color)
    if border_color:
        bc = border_color if isinstance(border_color, RGBColor) else parse_color(border_color)
        shape.line.color.rgb = bc
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()

    # Apply corner radius if specified
    if corner_radius is not None:
        # Corner radius is set via the shape's adjustment values
        # The adjustment value is a proportion of the shorter dimension
        shorter = min(width, height)
        if shorter > 0:
            adj_val = int((corner_radius / shorter) * 100000)
            shape.adjustments[0] = min(adj_val / 100000.0, 0.5)

    # Determine effective font size
    fs = Pt(font_size) if font_size else Pt(22)

    if rich_text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        for segment in rich_text:
            run = p.add_run()
            run.text = segment.get("text", "")
            run.font.name = IBM_FONT
            run.font.size = fs
            if segment.get("bold"):
                run.font.bold = True
            if segment.get("color"):
                run.font.color.rgb = parse_color(segment["color"])
    elif text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = IBM_FONT
        run.font.size = fs
        if text_color:
            run.font.color.rgb = parse_color(text_color)

    # Apply vertical alignment
    if valign and (text or rich_text):
        tf = shape.text_frame
        valign_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
        if valign in valign_map:
            tf.paragraphs[0].alignment = None  # Reset any horizontal alignment
            # Set vertical anchor on the text frame body properties
            txBody = shape._element.txBody
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
                bodyPr.set('anchor', anchor_map[valign])

    # Apply shadow
    if shadow:
        _apply_shadow(shape)

    return shape


def _add_cards(slide, cards, base_dir):
    """Add compound card visual elements to the slide.

    Each card consists of a background rectangle, top accent bar, title text,
    body text, and optional icon overlay.

    Args:
        slide: The slide to add cards to.
        cards: List of card spec dicts.
        base_dir: Base directory for resolving image paths.
    """
    for card in cards:
        cx = card.get("x", 0)
        cy = card.get("y", 0)
        cw = card.get("width", 4.0)
        ch = card.get("height", 3.0)
        fill = card.get("fill", "#F4F4F4")
        accent_color = card.get("accent_color", "#0F62FE")
        title = card.get("title", "")
        body = card.get("body", "")
        icon = card.get("icon")
        shadow = card.get("shadow", False)
        border_color = card.get("border_color")
        corner_radius = card.get("corner_radius", 0.15)
        card_title_font_size = card.get("title_font_size", 28)
        card_body_font_size = card.get("body_font_size", 22)

        # 1. Background rectangle (rounded)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx), Inches(cy), Inches(cw), Inches(ch)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = parse_color(fill)
        if border_color:
            bg_shape.line.color.rgb = parse_color(border_color)
            bg_shape.line.width = Pt(1)
        else:
            bg_shape.line.fill.background()

        # Apply corner radius
        shorter = min(cw, ch)
        if shorter > 0:
            adj_val = int((corner_radius / shorter) * 100000)
            bg_shape.adjustments[0] = min(adj_val / 100000.0, 0.5)

        if shadow:
            _apply_shadow(bg_shape)

        # 2. Top accent bar (thin stripe at top of card)
        accent_height = 0.25
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx), Inches(cy), Inches(cw), Inches(accent_height)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = parse_color(accent_color)
        accent_shape.line.fill.background()

        # 3. Title text box (bold, positioned below accent bar)
        title_y = cy + accent_height + 0.2
        title_h = 0.8
        if title:
            txBox = slide.shapes.add_textbox(
                Inches(cx + 0.35), Inches(title_y),
                Inches(cw - 0.7), Inches(title_h)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.name = IBM_FONT
            run.font.size = Pt(card_title_font_size)
            run.font.bold = True
            run.font.color.rgb = parse_color("#161616")

        # 4. Body text box (regular, positioned below title)
        body_y = title_y + title_h + 0.05
        body_h = ch - (body_y - cy) - 0.2
        if body and body_h > 0:
            txBox = slide.shapes.add_textbox(
                Inches(cx + 0.35), Inches(body_y),
                Inches(cw - 0.7), Inches(body_h)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            if isinstance(body, list):
                # Render each list item as a bullet paragraph
                for i, item in enumerate(body):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"\u2022 {item}"
                    p.font.name = IBM_FONT
                    p.font.size = Pt(card_body_font_size)
                    p.font.color.rgb = parse_color("#525252")
            else:
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = body
                run.font.name = IBM_FONT
                run.font.size = Pt(card_body_font_size)
                run.font.color.rgb = parse_color("#525252")

        # 5. Optional icon overlay image
        if icon and base_dir:
            try:
                img_path = _resolve_image_path(icon, base_dir)
                icon_size = 1.0
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(cx + cw - icon_size - 0.2),
                    Inches(cy + accent_height + 0.1),
                    Inches(icon_size),
                    Inches(icon_size)
                )
            except Exception as e:
                print(f"  WARNING: Could not add card icon {icon}: {e}", file=sys.stderr)


def _add_text_boxes(slide, text_boxes):
    """Add freeform text boxes at specified positions on the slide.

    Args:
        slide: The slide to add text boxes to.
        text_boxes: List of text box spec dicts with position, size, and text properties.
    """
    for tb in text_boxes:
        tx = tb.get("x", 0)
        ty = tb.get("y", 0)
        tw = tb.get("width", 4.0)
        th = tb.get("height", 1.0)
        text = tb.get("text", "")
        font_size = tb.get("font_size", 22)
        bold = tb.get("bold", False)
        color = tb.get("color", "#161616")
        font = tb.get("font", IBM_FONT)
        align = tb.get("align", "left")
        valign = tb.get("valign")
        rich_text = tb.get("rich_text")

        txBox = slide.shapes.add_textbox(
            Inches(tx), Inches(ty), Inches(tw), Inches(th)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Set vertical alignment
        if valign:
            bodyPr = txBox._element.txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
                if valign in anchor_map:
                    bodyPr.set('anchor', anchor_map[valign])

        p = tf.paragraphs[0]

        # Set horizontal alignment
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                     "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
        if align in align_map:
            p.alignment = align_map[align]

        if rich_text:
            for segment in rich_text:
                run = p.add_run()
                run.text = segment.get("text", "")
                run.font.name = segment.get("font", font)
                run.font.size = Pt(segment.get("font_size", font_size))
                if segment.get("bold"):
                    run.font.bold = True
                if segment.get("color"):
                    run.font.color.rgb = parse_color(segment["color"])
        elif text:
            run = p.add_run()
            run.text = text
            run.font.name = font
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = parse_color(color)


def _style_chart(chart, font_name=IBM_FONT, show_legend=False):
    """Apply IBM brand styling to a chart."""
    chart.has_legend = show_legend
    if show_legend and chart.legend:
        chart.legend.font.name = font_name
        chart.legend.font.size = Pt(14)
        chart.legend.include_in_layout = False
    # Style axes if available (doughnut/pie charts don't have axes)
    try:
        chart.category_axis.has_title = False
        chart.category_axis.tick_labels.font.name = font_name
        chart.category_axis.tick_labels.font.size = Pt(14)
    except (ValueError, AttributeError):
        pass
    try:
        chart.value_axis.has_title = False
        chart.value_axis.tick_labels.font.name = font_name
        chart.value_axis.tick_labels.font.size = Pt(14)
    except (ValueError, AttributeError):
        pass


def _add_native_donut_chart(slide, spec):
    """Add a native PowerPoint donut chart with per-segment colors."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 4.0)
    h = spec.get("height", 4.0)

    # Force square dimensions
    side = min(w, h)
    w = h = side

    segments = spec.get("segments", [])
    if not segments:
        return

    chart_data = CategoryChartData()
    chart_data.categories = [s.get("label", "") for s in segments]
    chart_data.add_series("Data", [s.get("value", 0) for s in segments])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    )
    chart = chart_frame.chart

    # Apply per-segment colors
    plot = chart.plots[0]
    series = plot.series[0]
    for i, seg in enumerate(segments):
        if seg.get("color"):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = parse_color(seg["color"])

    # Set hole size via XML manipulation
    hole_pct = int(spec.get("hole_size", 0.5) * 100)
    doughnut_chart = chart.plots[0]._element
    nsmap = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
    existing = doughnut_chart.findall(qn('c:holeSize'))
    for e in existing:
        doughnut_chart.remove(e)
    hole_elem = etree.SubElement(doughnut_chart, qn('c:holeSize'))
    hole_elem.set('val', str(hole_pct))

    # Center label overlay
    if spec.get("center_label"):
        label_text = spec["center_label"]
        label_w = side * 0.5
        label_h = side * 0.3
        label_x = x + (side - label_w) / 2
        label_y = y + (side - label_h) / 2
        txBox = slide.shapes.add_textbox(
            Inches(label_x), Inches(label_y),
            Inches(label_w), Inches(label_h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(label_text)
        run.font.name = IBM_FONT
        run.font.bold = True
        run.font.size = Pt(28)
        run.font.color.rgb = parse_color(spec.get("center_label_color", "#161616"))
        bodyPr = txBox._element.txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')

    _style_chart(chart, show_legend=spec.get("show_legend", True))


def _add_native_bar_chart(slide, spec):
    """Add a native PowerPoint horizontal bar chart."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 6.0)
    h = spec.get("height", 4.0)

    data = spec.get("data", [])
    if not data:
        return

    chart_data = CategoryChartData()
    chart_data.categories = [d.get("label", "") for d in data]
    chart_data.add_series("Values", [d.get("value", 0) for d in data])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    )
    chart = chart_frame.chart

    # Per-bar colors
    plot = chart.plots[0]
    series = plot.series[0]
    for i, d in enumerate(data):
        if d.get("color"):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = parse_color(d["color"])

    # Show value data labels
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.font.name = IBM_FONT
    data_labels.font.size = Pt(14)
    data_labels.show_value = True

    # Set gap width
    plot.gap_width = 80

    _style_chart(chart, show_legend=False)


def _add_native_comparison_bars(slide, spec):
    """Add a native PowerPoint comparison bar chart with Before/After series."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 6.0)
    h = spec.get("height", 4.0)

    items = spec.get("items", [])
    if not items:
        return

    before_label = spec.get("before_label", "Before")
    after_label = spec.get("after_label", "After")
    before_color = spec.get("before_color", "#E0E0E0")
    after_color = spec.get("after_color", "#0F62FE")

    chart_data = CategoryChartData()
    chart_data.categories = [item.get("label", "") for item in items]
    chart_data.add_series(before_label, [item.get("before", 0) for item in items])
    chart_data.add_series(after_label, [item.get("after", 0) for item in items])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    )
    chart = chart_frame.chart

    # Color the two series
    plot = chart.plots[0]
    before_series = plot.series[0]
    before_series.format.fill.solid()
    before_series.format.fill.fore_color.rgb = parse_color(before_color)

    after_series = plot.series[1]
    after_series.format.fill.solid()
    after_series.format.fill.fore_color.rgb = parse_color(after_color)

    _style_chart(chart, show_legend=True)
    if chart.legend:
        chart.legend.font.name = IBM_FONT
        chart.legend.font.size = Pt(14)


def _add_native_line_chart(slide, spec):
    """Add a native PowerPoint line/sparkline chart."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 4.0)
    h = spec.get("height", 2.0)

    values = spec.get("values", [])
    if not values:
        return

    chart_data = CategoryChartData()
    chart_data.categories = list(range(1, len(values) + 1))
    chart_data.add_series("Trend", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    )
    chart = chart_frame.chart

    # Style the line
    series = chart.plots[0].series[0]
    series.format.line.color.rgb = parse_color(spec.get("color", "#0F62FE"))
    series.format.line.width = Pt(2.5)

    # Minimal chrome: hide axes, no legend, no gridlines
    chart.has_legend = False
    chart.category_axis.visible = False
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False


def _add_native_progress_ring(slide, spec):
    """Add a native PowerPoint progress ring (donut with 2 segments)."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 3.0)
    h = spec.get("height", 3.0)

    # Force square
    side = min(w, h)
    w = h = side

    value = spec.get("value", 75)
    max_val = spec.get("max_value", 100)
    color = spec.get("color", "#0F62FE")
    track_color = spec.get("track_color", "#E0E0E0")

    chart_data = CategoryChartData()
    chart_data.categories = ["Value", "Remaining"]
    chart_data.add_series("Progress", [value, max_val - value])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    )
    chart = chart_frame.chart

    # Color the segments
    plot = chart.plots[0]
    series = plot.series[0]
    series.points[0].format.fill.solid()
    series.points[0].format.fill.fore_color.rgb = parse_color(color)
    series.points[1].format.fill.solid()
    series.points[1].format.fill.fore_color.rgb = parse_color(track_color)

    # Large hole size (70%)
    doughnut_chart = plot._element
    existing = doughnut_chart.findall(qn('c:holeSize'))
    for e in existing:
        doughnut_chart.remove(e)
    hole_elem = etree.SubElement(doughnut_chart, qn('c:holeSize'))
    hole_elem.set('val', '70')

    # Hide legend and labels
    chart.has_legend = False

    # Overlay centered textbox with percentage
    label = spec.get("label", f"{value}%")
    label_w = side * 0.5
    label_h = side * 0.4
    label_x = x + (side - label_w) / 2
    label_y = y + (side - label_h) / 2
    txBox = slide.shapes.add_textbox(
        Inches(label_x), Inches(label_y),
        Inches(label_w), Inches(label_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(label)
    run.font.name = IBM_FONT
    run.font.bold = True
    run.font.size = Pt(48)
    run.font.color.rgb = parse_color(color)
    bodyPr = txBox._element.txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def _add_native_process_flow(slide, spec):
    """Add a native PowerPoint process flow with connected step shapes."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    total_width = spec.get("width", 20.0)
    h = spec.get("height", 2.5)

    steps = spec.get("steps", [])
    if not steps:
        return

    n = len(steps)
    arrow_w = 0.8
    step_w = (total_width - arrow_w * (n - 1)) / n

    for i, step in enumerate(steps):
        step_x = x + i * (step_w + arrow_w)
        step_color = step.get("color", "#0F62FE")

        # Draw rounded rectangle for the step
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(step_x), Inches(y), Inches(step_w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(step_color)
        shape.line.fill.background()

        # Add text to step shape
        tf = shape.text_frame
        tf.word_wrap = True
        bodyPr = shape._element.txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step.get("label", "")
        run.font.name = IBM_FONT
        run.font.bold = True
        run.font.size = Pt(22)
        run.font.color.rgb = parse_color("#FFFFFF")

        # Optional sublabel
        sublabel = step.get("sublabel", "")
        if sublabel:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = sublabel
            run2.font.name = IBM_FONT
            run2.font.size = Pt(18)
            run2.font.color.rgb = parse_color("#E0E0E0")

        # Draw arrow between steps (not after last step)
        if i < n - 1:
            arrow_x = step_x + step_w + 0.05
            arrow_y_center = y + h / 2
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arrow_x), Inches(arrow_y_center - 0.35),
                Inches(arrow_w - 0.1), Inches(0.7)
            )
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = parse_color("#A8A8A8")
            arrow_shape.line.fill.background()


def _add_native_stat_card(slide, spec):
    """Add a native PowerPoint stat card with accent bar and value display."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 4.0)
    h = spec.get("height", 3.5)
    accent_color = spec.get("accent_color", "#0F62FE")
    value = spec.get("value", "")
    label = spec.get("label", "")
    trend = spec.get("trend")

    # 1. Background rounded rectangle (white fill, gray border)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = parse_color("#FFFFFF")
    bg_shape.line.color.rgb = parse_color("#E0E0E0")
    bg_shape.line.width = Pt(1)
    _apply_shadow(bg_shape)

    # 2. Accent bar at top
    accent_h = 0.25
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(accent_h)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = parse_color(accent_color)
    accent_shape.line.fill.background()

    # 3. Large value text (centered)
    value_y = y + accent_h + 0.3
    value_h = 1.2
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(value_y),
        Inches(w - 0.4), Inches(value_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(value)
    run.font.name = IBM_FONT
    run.font.bold = True
    run.font.size = Pt(48)
    run.font.color.rgb = parse_color(accent_color)

    # 4. Label text (centered, gray)
    label_y = value_y + value_h + 0.1
    label_h = 0.8
    txBox2 = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(label_y),
        Inches(w - 0.4), Inches(label_h)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = str(label)
    run2.font.name = IBM_FONT
    run2.font.size = Pt(22)
    run2.font.color.rgb = parse_color("#525252")

    # 5. Optional trend text with arrow
    # trend can be a string ("up"/"down") or a dict {"text", "direction", "color"}
    trend_value = spec.get("trend_value", "")
    if trend:
        trend_y = label_y + label_h + 0.05
        trend_h = 0.6
        if isinstance(trend, dict):
            trend_text = trend.get("text", "")
            trend_dir = trend.get("direction", "up")
            trend_color = trend.get("color", "#24A148" if trend_dir == "up" else "#DA1E28")
        else:
            trend_text = str(trend_value)
            trend_dir = str(trend)
            trend_color = "#24A148" if trend_dir == "up" else "#DA1E28"
        arrow_char = "\u2191" if trend_dir == "up" else "\u2193"
        txBox3 = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(trend_y),
            Inches(w - 0.4), Inches(trend_h)
        )
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = f"{arrow_char} {trend_text}"
        run3.font.name = IBM_FONT
        run3.font.size = Pt(18)
        run3.font.color.rgb = parse_color(trend_color)


def _add_native_metric_card(slide, spec):
    """Add a simple native PowerPoint metric card."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 3.5)
    h = spec.get("height", 2.5)
    value = spec.get("value", "")
    label = spec.get("label", "")
    color = spec.get("color", "#0F62FE")

    # Rounded rectangle with subtle border
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = parse_color("#FFFFFF")
    bg_shape.line.color.rgb = parse_color("#E0E0E0")
    bg_shape.line.width = Pt(1)

    # Large centered value text
    value_h = h * 0.55
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(y + 0.2),
        Inches(w - 0.4), Inches(value_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(value)
    run.font.name = IBM_FONT
    run.font.bold = True
    run.font.size = Pt(48)
    run.font.color.rgb = parse_color(color)
    bodyPr = txBox._element.txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'b')

    # Label below
    label_y = y + 0.2 + value_h + 0.05
    label_h = h - value_h - 0.45
    txBox2 = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(label_y),
        Inches(w - 0.4), Inches(max(label_h, 0.5))
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = str(label)
    run2.font.name = IBM_FONT
    run2.font.size = Pt(22)
    run2.font.color.rgb = parse_color("#525252")


def _add_native_icon_badge(slide, spec):
    """Add a native PowerPoint oval icon badge with centered text."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 1.5)
    h = spec.get("height", 1.5)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_color(spec.get("color", "#0F62FE"))
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(spec.get("number", ""))
    run.font.name = IBM_FONT
    run.font.bold = True
    run.font.color.rgb = parse_color(spec.get("text_color", "#FFFFFF"))
    run.font.size = Pt(int(min(w, h) * 72 * 0.35))  # Scale to badge size

    # Center vertically
    bodyPr = shape._element.txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


def _add_native_quote_mark(slide, spec):
    """Add a large decorative quote mark."""
    x = spec.get("x", 0)
    y = spec.get("y", 0)
    w = spec.get("width", 2.0)
    h = spec.get("height", 2.0)
    color = spec.get("color", "#0F62FE")

    style = spec.get("style", "open")
    char = "\u201C" if style == "open" else "\u201D"
    font_size = int(w * 72 * 0.7)  # Scale to box size

    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = char
    run.font.name = IBM_FONT
    run.font.size = Pt(font_size)
    run.font.color.rgb = parse_color(color)


def _render_png_visual(slide, vis, base_dir):
    """Render a visualization as PNG and embed it."""
    try:
        from .svg_visuals import render_visual
    except ImportError:
        try:
            svg_visuals_path = Path(__file__).parent / "svg_visuals.py"
            if svg_visuals_path.exists():
                import importlib.util
                mod_spec = importlib.util.spec_from_file_location("svg_visuals", svg_visuals_path)
                svg_visuals = importlib.util.module_from_spec(mod_spec)
                mod_spec.loader.exec_module(svg_visuals)
                render_visual = svg_visuals.render_visual
            else:
                print(f"  WARNING: svg_visuals.py not found, skipping visual {vis.get('type')}", file=sys.stderr)
                return
        except Exception as e:
            print(f"  WARNING: Could not import svg_visuals: {e}", file=sys.stderr)
            return

    try:
        png_bytes = render_visual(vis)
        if png_bytes:
            import io
            img_stream = io.BytesIO(png_bytes)
            slide.shapes.add_picture(
                img_stream,
                Inches(vis.get("x", 0)),
                Inches(vis.get("y", 0)),
                Inches(vis.get("width", 2)),
                Inches(vis.get("height", 2))
            )
    except Exception as e:
        print(f"  WARNING: Could not render visual {vis.get('type')}: {e}", file=sys.stderr)


def _add_native_architecture_diagram(slide, spec):
    """Add a layered architecture diagram using native PowerPoint shapes.

    Each layer is a wide rounded rectangle with a colored accent bar on the left,
    a label, and horizontally arranged item shapes.  Thin connector lines join
    adjacent layers.
    """
    x = spec.get("x", 0.63)
    y = spec.get("y", 2.0)
    width = spec.get("width", 25.4)
    height = spec.get("height", 11.0)
    layers = spec.get("layers", [])
    draw_connectors = spec.get("connectors", True)

    if not layers:
        return

    n_layers = len(layers)
    layer_gap = 0.3  # gap between layers in inches
    connector_h = 0.2  # height of connectors in inches
    available_h = height - (n_layers - 1) * (layer_gap + connector_h)
    layer_h = available_h / n_layers

    accent_w = 0.4  # accent bar width
    label_w = 3.0  # label area width
    label_offset = 0.5  # label x offset from layer left
    items_x_offset = 4.0  # where items start (from layer left)
    item_margin = 0.2  # gap between item shapes
    item_v_pad = 0.25  # vertical padding inside layer for items

    for i, layer in enumerate(layers):
        ly = y + i * (layer_h + layer_gap + connector_h)
        layer_color_str = layer.get("color", "#0F62FE")
        layer_color = parse_color(layer_color_str)
        label_text = layer.get("label", "")
        items = layer.get("items", [])

        # -- Layer background (white fill, subtle border) --
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(ly), Inches(width), Inches(layer_h),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bg_shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        bg_shape.line.width = Pt(1)
        # No text on the background shape
        bg_shape.text_frame.clear()

        # -- Accent bar (left edge) --
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(ly), Inches(accent_w), Inches(layer_h),
        )
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = layer_color
        bar_shape.line.fill.background()
        bar_shape.text_frame.clear()

        # -- Layer label --
        label_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + label_offset), Inches(ly), Inches(label_w), Inches(layer_h),
        )
        label_shape.fill.background()
        label_shape.line.fill.background()
        tf = label_shape.text_frame
        tf.word_wrap = True
        bodyPr = label_shape._element.txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label_text
        run.font.name = IBM_FONT
        run.font.bold = True
        run.font.size = Pt(24)
        run.font.color.rgb = layer_color

        # -- Item shapes --
        if items:
            items_area_w = width - items_x_offset - 0.3  # leave right margin
            n_items = len(items)
            item_w = (items_area_w - item_margin * (n_items - 1)) / n_items
            item_h = layer_h - 2 * item_v_pad

            for j, item_text in enumerate(items):
                ix = x + items_x_offset + j * (item_w + item_margin)
                iy = ly + item_v_pad

                item_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(ix), Inches(iy), Inches(item_w), Inches(item_h),
                )
                item_shape.fill.solid()
                item_shape.fill.fore_color.rgb = layer_color
                item_shape.line.fill.background()

                # Shadow on item shapes
                _apply_shadow(item_shape)

                itf = item_shape.text_frame
                itf.word_wrap = True
                item_bodyPr = item_shape._element.txBody.find(qn('a:bodyPr'))
                if item_bodyPr is not None:
                    item_bodyPr.set('anchor', 'ctr')
                ip = itf.paragraphs[0]
                ip.alignment = PP_ALIGN.CENTER
                irun = ip.add_run()
                irun.text = item_text
                irun.font.name = IBM_FONT
                irun.font.size = Pt(18)
                irun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # -- Connector line to next layer --
        if draw_connectors and i < n_layers - 1:
            conn_y = ly + layer_h
            conn_x = x + width / 2 - 0.05  # center horizontally
            conn_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(conn_x), Inches(conn_y),
                Inches(0.1), Inches(layer_gap + connector_h),
            )
            conn_shape.fill.solid()
            conn_shape.fill.fore_color.rgb = RGBColor(0xA8, 0xA8, 0xA8)
            conn_shape.line.fill.background()
            conn_shape.text_frame.clear()


def _add_visuals(slide, visuals, base_dir):
    """Render and place data visualizations on the slide.

    By default uses native PowerPoint shapes/charts for editable output.
    Set "render_mode": "png" on a visual for legacy bitmap rendering.
    """
    NATIVE_RENDERERS = {
        "donut_chart": _add_native_donut_chart,
        "horizontal_bars": _add_native_bar_chart,
        "comparison_bars": _add_native_comparison_bars,
        "sparkline": _add_native_line_chart,
        "progress_ring": _add_native_progress_ring,
        "process_flow": _add_native_process_flow,
        "stat_card": _add_native_stat_card,
        "metric_card": _add_native_metric_card,
        "icon_badge": _add_native_icon_badge,
        "quote_mark": _add_native_quote_mark,
        "architecture_diagram": _add_native_architecture_diagram,
    }
    PNG_ONLY_TYPES = {"gradient_bar", "accent_gradient"}

    for vis in visuals:
        vis_type = vis.get("type", "")
        render_mode = vis.get("render_mode", "native")

        if render_mode == "native" and vis_type in NATIVE_RENDERERS:
            try:
                NATIVE_RENDERERS[vis_type](slide, vis)
            except Exception as e:
                print(f"  WARNING: Native render failed for {vis_type}, falling back to PNG: {e}", file=sys.stderr)
                _render_png_visual(slide, vis, base_dir)
        else:
            _render_png_visual(slide, vis, base_dir)


# ---------------------------------------------------------------------------
# Text formatting helpers
# ---------------------------------------------------------------------------

def _set_text_with_font(
    text_frame,
    text: str,
    font_name: str = IBM_FONT,
    font_size: Optional[Pt] = None,
    font_color: Optional[RGBColor] = None,
    bold: Optional[bool] = None,
    alignment: Optional[PP_ALIGN] = None,
):
    """Set the text of a text frame, clearing existing content and applying formatting."""
    # Snapshot inherited paragraph properties before clearing
    try:
        first_pPr = text_frame.paragraphs[0]._p.find(qn('a:pPr'))
        inherited_pPr = copy.deepcopy(first_pPr) if first_pPr is not None else None
    except Exception:
        inherited_pPr = None

    text_frame.clear()
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]

    # Restore inherited paragraph properties
    if inherited_pPr is not None:
        try:
            p._p.insert(0, inherited_pPr)
        except Exception:
            pass

    if alignment is not None:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    if font_size is not None:
        run.font.size = font_size
    if font_color is not None:
        run.font.color.rgb = font_color
    if bold is not None:
        run.font.bold = bold


def _set_bullets(
    text_frame,
    items: list[str],
    font_name: str = IBM_FONT,
    font_size: Optional[Pt] = None,
    font_color: Optional[RGBColor] = None,
):
    """Write a list of strings as bullet-point paragraphs."""
    text_frame.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.level = 0

        # Add bullet character and spacing
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u2022')
        # Add space after each bullet paragraph
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts.set('val', '600')  # 6pt spacing after

        run = p.add_run()
        run.text = item
        run.font.name = font_name
        if font_size is not None:
            run.font.size = font_size
        if font_color is not None:
            run.font.color.rgb = font_color


def _fill_text_placeholder(placeholder, content, font_color: Optional[RGBColor] = None):
    """Fill a text placeholder with either a string or a list of bullet strings."""
    if content is None:
        return
    tf = placeholder.text_frame
    if isinstance(content, list):
        _set_bullets(tf, content, font_color=font_color)
    else:
        _set_text_with_font(tf, str(content), font_color=font_color)


def _resolve_image_path(image_path: str, base_dir: Optional[Path] = None) -> Path:
    """Resolve an image path, trying relative to base_dir first, then absolute."""
    p = Path(image_path)
    if p.is_absolute() and p.exists():
        return p
    if base_dir is not None:
        candidate = base_dir / image_path
        if candidate.exists():
            return candidate
    if p.exists():
        return p.resolve()
    raise FileNotFoundError(f"Image not found: {image_path}")


def _insert_image_placeholder(placeholder, image_path: str, base_dir: Optional[Path] = None):
    """Insert an image into a picture placeholder.

    SVG files cannot be inserted into picture placeholders (python-pptx limitation).
    If an SVG is provided, try to find a PNG fallback with the same base name.
    """
    resolved = _resolve_image_path(image_path, base_dir)
    if resolved.suffix.lower() == ".svg":
        # Try PNG fallback: same directory, replace .svg with .png
        # Also try the paired naming convention from the IBM template:
        # e.g., pptx_image90.svg -> pptx_image89.png (svg at N, png at N-1)
        png_candidate = resolved.with_suffix(".png")
        if png_candidate.exists():
            resolved = png_candidate
        else:
            # Try finding a PNG with a nearby numeric name (fallback pairs)
            stem = resolved.stem
            import re
            m = re.match(r"(.+?)(\d+)$", stem)
            if m:
                prefix, num = m.group(1), int(m.group(2))
                for offset in [-1, 1, -2, 2]:
                    candidate = resolved.parent / f"{prefix}{num + offset}.png"
                    if candidate.exists():
                        resolved = candidate
                        break
                else:
                    print(f"  WARNING: SVG not supported and no PNG fallback found for {image_path}",
                          file=sys.stderr)
                    return
            else:
                print(f"  WARNING: SVG not supported and no PNG fallback found for {image_path}",
                      file=sys.stderr)
                return
    try:
        placeholder.insert_picture(str(resolved))
    except Exception as e:
        print(f"  WARNING: Could not insert image {resolved}: {e}", file=sys.stderr)


# ---------------------------------------------------------------------------
# Slide builders by category
# ---------------------------------------------------------------------------

def _build_cover_slide(slide, spec: dict, meta: dict, base_dir: Optional[Path] = None,
                       text_color: Optional[RGBColor] = None):
    """Build cover layouts (0-5)."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    for idx, role in ph_map.items():
        ph = placeholders.get(idx)
        if ph is None:
            continue
        if role == "title" and spec.get("title"):
            _fill_text_placeholder(ph, spec["title"], font_color=text_color)
        elif role == "subtitle" and spec.get("subtitle"):
            _fill_text_placeholder(ph, spec["subtitle"], font_color=text_color)
        elif role == "body" and spec.get("body"):
            _fill_text_placeholder(ph, spec["body"], font_color=text_color)
        elif role == "image" and spec.get("image"):
            _insert_image_placeholder(ph, spec["image"], base_dir)


def _build_navigation_slide(slide, spec: dict, meta: dict, base_dir: Optional[Path] = None,
                            text_color: Optional[RGBColor] = None):
    """Build navigation layouts (6-8): Contents, Section divider, Large text."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    for idx, role in ph_map.items():
        ph = placeholders.get(idx)
        if ph is None:
            continue
        if role == "title" and spec.get("title"):
            _fill_text_placeholder(ph, spec["title"], font_color=text_color)
        elif role == "body" and spec.get("body"):
            _fill_text_placeholder(ph, spec["body"], font_color=text_color)
        elif role == "body_right" and spec.get("body_right"):
            _fill_text_placeholder(ph, spec["body_right"], font_color=text_color)
        elif role == "footer" and spec.get("footer"):
            _fill_text_placeholder(ph, spec["footer"])


def _build_callout_slide(slide, spec: dict, meta: dict, base_dir: Optional[Path] = None,
                         text_color: Optional[RGBColor] = None):
    """Build callout layouts (9-10)."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    for idx, role in ph_map.items():
        ph = placeholders.get(idx)
        if ph is None:
            continue
        if role == "title" and spec.get("title"):
            _fill_text_placeholder(ph, spec["title"], font_color=text_color)
        elif role == "body" and spec.get("body"):
            content = spec["body"]
            # For callout body, join list items with newlines if it's a list
            if isinstance(content, list):
                content = "\n".join(content)
            _fill_text_placeholder(ph, content, font_color=text_color)


def _build_data_slide(slide, spec: dict, meta: dict, base_dir: Optional[Path] = None,
                      text_color: Optional[RGBColor] = None):
    """Build data layouts (11-14) with data_points."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    data_points = spec.get("data_points", [])

    for idx, role in ph_map.items():
        ph = placeholders.get(idx)
        if ph is None:
            continue

        if role == "title" and spec.get("title"):
            _fill_text_placeholder(ph, spec["title"], font_color=text_color)
        elif role == "body" and spec.get("body"):
            _fill_text_placeholder(ph, spec["body"], font_color=text_color)
        elif role == "footer" and spec.get("footer"):
            _fill_text_placeholder(ph, spec["footer"])
        elif role.startswith("data_value_"):
            dp_idx = int(role.split("_")[-1]) - 1
            if dp_idx < len(data_points):
                value = data_points[dp_idx].get("value", "")
                _set_text_with_font(
                    ph.text_frame, str(value),
                    font_name=IBM_FONT,
                    font_size=Pt(60),
                    font_color=text_color or IBM_BLUE_60,
                    bold=True,
                )
        elif role.startswith("data_label_"):
            dp_idx = int(role.split("_")[-1]) - 1
            if dp_idx < len(data_points):
                label = data_points[dp_idx].get("label", "")
                _fill_text_placeholder(ph, label, font_color=text_color)


def _build_text_columns_slide(slide, spec: dict, meta: dict, base_dir: Optional[Path] = None,
                              text_color: Optional[RGBColor] = None):
    """Build text column layouts (15-25)."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # Gather column content from spec: body is col_1, body_right is col_2, etc.
    # Also support explicit "columns" list in spec
    columns = spec.get("columns", [])
    body_items = spec.get("body", [])
    body_right_items = spec.get("body_right", [])
    icons = spec.get("icons", [])
    col_heads = spec.get("column_headings", [])

    for idx, role in ph_map.items():
        ph = placeholders.get(idx)
        if ph is None:
            continue

        if role == "title" and spec.get("title"):
            _fill_text_placeholder(ph, spec["title"], font_color=text_color)
        elif role == "subtitle" and spec.get("subtitle"):
            _fill_text_placeholder(ph, spec["subtitle"], font_color=text_color)
        elif role == "body" and body_items:
            _fill_text_placeholder(ph, body_items, font_color=text_color)
        elif role == "body_right" and body_right_items:
            _fill_text_placeholder(ph, body_right_items, font_color=text_color)
        elif role.startswith("col_") and not role.endswith("_head"):
            col_num = int(role.split("_")[1]) - 1
            if col_num < len(columns):
                _fill_text_placeholder(ph, columns[col_num], font_color=text_color)
            elif col_num == 0 and body_items:
                _fill_text_placeholder(ph, body_items, font_color=text_color)
            elif col_num == 1 and body_right_items:
                _fill_text_placeholder(ph, body_right_items, font_color=text_color)
        elif role.endswith("_head"):
            head_num = int(role.split("_")[1]) - 1
            if head_num < len(col_heads):
                _fill_text_placeholder(ph, col_heads[head_num], font_color=text_color)
        elif role.startswith("icon_"):
            icon_num = int(role.split("_")[1]) - 1
            if icon_num < len(icons):
                _insert_image_placeholder(ph, icons[icon_num], base_dir)
        elif role == "footer" and spec.get("footer"):
            _fill_text_placeholder(ph, spec["footer"])


def _build_box_grid_slide(slide, spec: dict, meta: dict, base_dir: Optional[Path] = None,
                          text_color: Optional[RGBColor] = None):
    """Build box grid layouts (26-34)."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # Box items from spec: use "boxes" list or fall back to "body" list
    boxes = spec.get("boxes", [])
    box_heads = spec.get("box_headings", [])
    icons = spec.get("icons", [])
    body_items = spec.get("body", [])

    for idx, role in ph_map.items():
        ph = placeholders.get(idx)
        if ph is None:
            continue

        if role == "title" and spec.get("title"):
            _fill_text_placeholder(ph, spec["title"], font_color=text_color)
        elif role == "subtitle" and spec.get("subtitle"):
            _fill_text_placeholder(ph, spec["subtitle"], font_color=text_color)
        elif role.startswith("box_") and not role.endswith("_head"):
            box_num = int(role.split("_")[1]) - 1
            if box_num < len(boxes):
                _fill_text_placeholder(ph, boxes[box_num], font_color=text_color)
            elif box_num < len(body_items):
                _fill_text_placeholder(ph, body_items[box_num], font_color=text_color)
        elif role.endswith("_head"):
            head_num = int(role.split("_")[1]) - 1
            if head_num < len(box_heads):
                _fill_text_placeholder(ph, box_heads[head_num], font_color=text_color)
        elif role.startswith("icon_"):
            icon_num = int(role.split("_")[1]) - 1
            if icon_num < len(icons):
                _insert_image_placeholder(ph, icons[icon_num], base_dir)
        elif role == "footer" and spec.get("footer"):
            _fill_text_placeholder(ph, spec["footer"])


def _build_media_slide(slide, spec: dict, meta: dict, base_dir: Optional[Path] = None,
                       text_color: Optional[RGBColor] = None):
    """Build media layouts (35-39)."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    for idx, role in ph_map.items():
        ph = placeholders.get(idx)
        if ph is None:
            continue

        if role == "title" and spec.get("title"):
            _fill_text_placeholder(ph, spec["title"], font_color=text_color)
        elif role == "subtitle" and spec.get("subtitle"):
            _fill_text_placeholder(ph, spec["subtitle"], font_color=text_color)
        elif role == "body" and spec.get("body"):
            _fill_text_placeholder(ph, spec["body"], font_color=text_color)
        elif role == "image" and spec.get("image"):
            _insert_image_placeholder(ph, spec["image"], base_dir)
        elif role == "footer" and spec.get("footer"):
            _fill_text_placeholder(ph, spec["footer"])


def _build_special_slide(slide, spec: dict, meta: dict, layout_idx: int,
                         base_dir: Optional[Path] = None,
                         text_color: Optional[RGBColor] = None):
    """Build special layouts (40-48)."""
    ph_map = meta["placeholders"]
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    if layout_idx == 40:
        # Contacts / profiles
        profiles = spec.get("profiles", [])
        for idx, role in ph_map.items():
            ph = placeholders.get(idx)
            if ph is None:
                continue
            if role == "title" and spec.get("title"):
                _fill_text_placeholder(ph, spec["title"], font_color=text_color)
            elif role.startswith("profile_image_"):
                prof_num = int(role.split("_")[-1]) - 1
                if prof_num < len(profiles) and profiles[prof_num].get("image"):
                    _insert_image_placeholder(ph, profiles[prof_num]["image"], base_dir)
            elif role.startswith("profile_info_"):
                prof_num = int(role.split("_")[-1]) - 1
                if prof_num < len(profiles):
                    profile = profiles[prof_num]
                    name = profile.get("name", "")
                    role_text = profile.get("role", "")
                    info_text = f"{name}\n{role_text}" if role_text else name
                    _fill_text_placeholder(ph, info_text, font_color=text_color)
            elif role == "footer" and spec.get("footer"):
                _fill_text_placeholder(ph, spec["footer"])

    elif layout_idx == 41:
        # Table -- fill title; table placeholder requires manual XML work
        # or using the table API. We fill what we can.
        for idx, role in ph_map.items():
            ph = placeholders.get(idx)
            if ph is None:
                continue
            if role == "title" and spec.get("title"):
                _fill_text_placeholder(ph, spec["title"], font_color=text_color)
            elif role == "table" and spec.get("table_data"):
                _fill_table_placeholder(ph, spec["table_data"], text_color=text_color)
            elif role == "footer" and spec.get("footer"):
                _fill_text_placeholder(ph, spec["footer"])

    elif layout_idx == 42:
        # Chart
        for idx, role in ph_map.items():
            ph = placeholders.get(idx)
            if ph is None:
                continue
            if role == "title" and spec.get("title"):
                _fill_text_placeholder(ph, spec["title"], font_color=text_color)
            elif role == "chart_title" and spec.get("subtitle"):
                _fill_text_placeholder(ph, spec["subtitle"], font_color=text_color)
            elif role == "chart_body" and spec.get("body"):
                _fill_text_placeholder(ph, spec["body"], font_color=text_color)
            elif role == "footer" and spec.get("footer"):
                _fill_text_placeholder(ph, spec["footer"])

    elif layout_idx in (43, 44):
        # Legal disclaimer
        for idx, role in ph_map.items():
            ph = placeholders.get(idx)
            if ph is None:
                continue
            if role == "legal_title" and spec.get("title"):
                _fill_text_placeholder(ph, spec["title"], font_color=text_color)
            elif role == "legal_body" and spec.get("body"):
                content = spec["body"]
                if isinstance(content, list):
                    content = "\n\n".join(content)
                _fill_text_placeholder(ph, content, font_color=text_color)
            elif role == "legal_body2" and spec.get("body_right"):
                content = spec["body_right"]
                if isinstance(content, list):
                    content = "\n\n".join(content)
                _fill_text_placeholder(ph, content, font_color=text_color)
            elif role == "footer" and spec.get("footer"):
                _fill_text_placeholder(ph, spec["footer"])

    elif layout_idx in (45, 46):
        # Blank slides
        for idx, role in ph_map.items():
            ph = placeholders.get(idx)
            if ph is None:
                continue
            if role == "title":
                if spec.get("title"):
                    _fill_text_placeholder(ph, spec["title"], font_color=text_color)
                else:
                    # Clear placeholder to prevent "Click to add title" overlay
                    ph.text_frame.clear()
            elif role == "footer":
                if spec.get("footer"):
                    _fill_text_placeholder(ph, spec["footer"])
                else:
                    ph.text_frame.clear()
            elif role == "slide_number":
                if not spec.get("slide_number"):
                    ph.text_frame.clear()

    # 47 (End slide) and 48 (DEFAULT) have no placeholders -- nothing to fill


def _fill_table_placeholder(placeholder, table_data: dict, text_color=None):
    """Fill a table placeholder with row/column data."""
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers and not rows:
        return

    num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    num_rows = (1 if headers else 0) + len(rows)
    if num_cols == 0 or num_rows == 0:
        return

    table_shape = placeholder.insert_table(num_rows, num_cols)
    table = table_shape.table

    row_offset = 0
    if headers:
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            # Dark header background for contrast
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x16)  # gray100
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = IBM_FONT
                    run.font.bold = True
                    run.font.size = Pt(22)
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # white text
        row_offset = 1

    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            if col_idx >= num_cols:
                break
            cell = table.cell(row_idx + row_offset, col_idx)
            cell.text = str(val)
            # Alternating row backgrounds
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # white
            else:
                cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF4)  # gray10
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = IBM_FONT
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(0x16, 0x16, 0x16)  # dark text for contrast
                    if text_color:
                        run.font.color.rgb = text_color


# ---------------------------------------------------------------------------
# Main slide builder dispatch
# ---------------------------------------------------------------------------

CATEGORY_BUILDERS = {
    "cover": _build_cover_slide,
    "navigation": _build_navigation_slide,
    "callout": _build_callout_slide,
    "data": _build_data_slide,
    "text_columns": _build_text_columns_slide,
    "box_grid": _build_box_grid_slide,
    "media": _build_media_slide,
}


def _add_programmatic_table(slide, table_spec):
    """Add a properly positioned and styled table directly on the slide.

    Unlike _fill_table_placeholder which uses template placeholders,
    this creates a table at exact coordinates for reliable positioning.

    table_spec: {
        "x": float, "y": float, "width": float, "height": float,
        "headers": ["Col1", "Col2", ...],
        "rows": [["val1", "val2", ...], ...],
        "col_widths": [float, ...] (optional, in inches)
    }
    """
    x = table_spec.get("x", 0.63)
    y = table_spec.get("y", 2.5)
    w = table_spec.get("width", 25.4)
    h = table_spec.get("height", 9.0)
    headers = table_spec.get("headers", [])
    rows = table_spec.get("rows", [])
    if not headers and not rows:
        return

    num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    num_rows = (1 if headers else 0) + len(rows)
    if num_cols == 0 or num_rows == 0:
        return

    # Create table at exact position
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = table_shape.table

    # Set column widths
    col_widths = table_spec.get("col_widths")
    if col_widths:
        for i, cw in enumerate(col_widths):
            if i < num_cols:
                table.columns[i].width = Inches(cw)
    else:
        # Equal column widths
        col_w = w / num_cols
        for i in range(num_cols):
            table.columns[i].width = Inches(col_w)

    # Header row
    row_offset = 0
    if headers:
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x16)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = IBM_FONT
                    run.font.bold = True
                    run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        row_offset = 1

    # Data rows with alternating backgrounds
    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            if col_idx >= num_cols:
                break
            cell = table.cell(row_idx + row_offset, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF4)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = IBM_FONT
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(0x16, 0x16, 0x16)
                    # Color checkmarks green, X marks red
                    if val.strip() in ("\u2713", "✓"):
                        run.font.color.rgb = RGBColor(0x19, 0x80, 0x38)
                    elif val.strip() in ("\u2717", "✗", "✕"):
                        run.font.color.rgb = RGBColor(0xDA, 0x1E, 0x28)


def _clear_unused_placeholders(slide):
    """Remove placeholder shapes that have no useful content.

    Removes the placeholder XML element entirely so no empty box
    appears on the slide. Handles both text and image placeholders.
    """
    # Collect elements to remove (can't modify while iterating)
    to_remove = []
    for ph in slide.placeholders:
        should_remove = False
        if hasattr(ph, 'text_frame'):
            text = ph.text_frame.text.strip() if ph.text_frame.text else ""
            if not text or text.startswith("Click to add") or text.startswith("Click to edit"):
                should_remove = True
        else:
            # Image placeholder with no image inserted — also remove
            should_remove = True
        if should_remove:
            to_remove.append(ph._element)
    # Remove collected elements from slide XML
    for sp in to_remove:
        sp.getparent().remove(sp)


def build_slide(prs: Presentation, slide_spec: dict, base_dir: Optional[Path] = None):
    """Add a single slide to the presentation based on the slide spec."""
    layout_idx = slide_spec.get("layout", 0)
    if layout_idx < 0 or layout_idx >= len(prs.slide_layouts):
        print(f"  WARNING: Layout index {layout_idx} out of range (0-{len(prs.slide_layouts)-1}), using 0",
              file=sys.stderr)
        layout_idx = 0

    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    meta = LAYOUT_META.get(layout_idx)
    if meta is None:
        print(f"  WARNING: No metadata for layout {layout_idx}, slide added but not populated",
              file=sys.stderr)
        return slide

    # Parse optional text_color override
    text_color = None
    if slide_spec.get("text_color"):
        text_color = parse_color(slide_spec["text_color"])

    category = meta["category"]
    if category == "special":
        _build_special_slide(slide, slide_spec, meta, layout_idx,
                             base_dir=base_dir, text_color=text_color)
    else:
        builder = CATEGORY_BUILDERS.get(category)
        if builder:
            builder(slide, slide_spec, meta, base_dir=base_dir, text_color=text_color)

    # Clear any unused placeholders to prevent "Click to add title" overlays
    _clear_unused_placeholders(slide)

    # --- Visual Enhancement Layer ---
    if slide_spec.get("background"):
        _set_slide_background(slide, slide_spec["background"])

    # Accent color for auto-generated accent bars
    accent_color = None
    if slide_spec.get("accent_color"):
        accent_color = parse_color(slide_spec["accent_color"])

    # Explicit accent bars
    if slide_spec.get("accent_bars"):
        for bar in slide_spec["accent_bars"]:
            _add_accent_bar(slide, bar["x"], bar["y"], bar["width"], bar["height"],
                           parse_color(bar["color"]))

    # Freeform image overlays
    if slide_spec.get("overlays"):
        _add_overlay_images(slide, slide_spec["overlays"], base_dir)

    # Dark scrim overlay for text contrast on image-backed slides
    if slide_spec.get("scrim"):
        scrim = slide_spec["scrim"]
        _add_scrim_overlay(
            slide,
            scrim.get("x", 0), scrim.get("y", 0),
            scrim.get("width", 26.67), scrim.get("height", 15.0),
            scrim.get("color", "#000000"),
            scrim.get("opacity", 50)
        )

    # Divider lines
    if slide_spec.get("dividers"):
        for d in slide_spec["dividers"]:
            _add_divider_line(slide, d["x"], d["y"], d["length"],
                            d.get("orientation", "vertical"),
                            parse_color(d["color"]) if d.get("color") else None)

    # Callout/highlight shapes
    if slide_spec.get("callouts"):
        for c in slide_spec["callouts"]:
            _add_callout_shape(
                slide, c["x"], c["y"], c["width"], c["height"],
                c["fill"], c.get("border"), c.get("text"),
                rich_text=c.get("rich_text"),
                shadow=c.get("shadow", False),
                font_size=c.get("font_size"),
                valign=c.get("valign"),
                text_color=c.get("text_color"),
                corner_radius=c.get("corner_radius"),
            )

    # Cards (compound visual elements)
    if slide_spec.get("cards"):
        _add_cards(slide, slide_spec["cards"], base_dir)

    # Freeform text boxes
    if slide_spec.get("text_boxes"):
        _add_text_boxes(slide, slide_spec["text_boxes"])

    # Programmatic table (positioned via spec, not template placeholder)
    if slide_spec.get("table"):
        _add_programmatic_table(slide, slide_spec["table"])

    # Data visualizations
    if slide_spec.get("visuals"):
        _add_visuals(slide, slide_spec["visuals"], base_dir)

    # Speaker notes
    if slide_spec.get("notes"):
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        _set_text_with_font(notes_tf, slide_spec["notes"], font_size=Pt(12))

    return slide


# ---------------------------------------------------------------------------
# Presentation-level metadata
# ---------------------------------------------------------------------------

def _set_core_properties(prs: Presentation, spec: dict):
    """Set presentation-level metadata from the spec."""
    core = prs.core_properties
    if spec.get("title"):
        core.title = spec["title"]
    if spec.get("author"):
        core.author = spec["author"]
    if spec.get("subtitle"):
        core.subject = spec["subtitle"]


# ---------------------------------------------------------------------------
# Main build function
# ---------------------------------------------------------------------------

def build_presentation(spec: dict, base_dir: Optional[Path] = None) -> str:
    """Build a complete presentation from a JSON specification.

    Args:
        spec: The parsed JSON specification dict.
        base_dir: Directory for resolving relative image paths.

    Returns:
        The output file path.
    """
    prs = load_template()
    _set_core_properties(prs, spec)

    # Remove all pre-existing slides from the template so we start clean.
    # The template may contain sample/demo slides that we don't want.
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    slides = spec.get("slides", [])
    if not slides:
        print("WARNING: No slides specified in the JSON spec.", file=sys.stderr)

    for i, slide_spec in enumerate(slides):
        layout_idx = slide_spec.get("layout", 0)
        layout_name = LAYOUT_META.get(layout_idx, {}).get("name", "Unknown")
        print(f"  Slide {i+1}: layout {layout_idx} ({layout_name})")
        build_slide(prs, slide_spec, base_dir=base_dir)

    output_file = spec.get("output_file", "presentation.pptx")
    prs.save(output_file)
    return output_file


# ---------------------------------------------------------------------------
# Layout listing
# ---------------------------------------------------------------------------

def list_layouts():
    """Print all available layouts with their placeholder details."""
    prs = load_template()
    print(f"IBM Presentation Template: {len(prs.slide_layouts)} layouts")
    print(f"Slide dimensions: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
    print("=" * 80)
    for i, layout in enumerate(prs.slide_layouts):
        meta = LAYOUT_META.get(i, {})
        category = meta.get("category", "unknown")
        print(f"\nLayout {i}: \"{layout.name}\"  [{category}]")
        ph_roles = meta.get("placeholders", {})
        for ph in layout.placeholders:
            idx = ph.placeholder_format.idx
            role = ph_roles.get(idx, "?")
            print(f"  idx={idx:>2d}  type={str(ph.placeholder_format.type):>25s}  role={role}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Build IBM-branded PowerPoint presentations from JSON specifications.",
        epilog="Example: python build_presentation.py spec.json -o output.pptx",
    )
    parser.add_argument(
        "input",
        nargs="?",
        default=None,
        help="JSON spec file path, or '-' for stdin. Omit to use --list-layouts.",
    )
    parser.add_argument(
        "-o", "--output",
        default=None,
        help="Output PPTX file path (overrides spec's output_file).",
    )
    parser.add_argument(
        "--list-layouts",
        action="store_true",
        help="Print all available slide layouts and exit.",
    )
    parser.add_argument(
        "--base-dir",
        default=None,
        help="Base directory for resolving relative image paths.",
    )

    args = parser.parse_args()

    if args.list_layouts:
        list_layouts()
        return

    if args.input is None:
        parser.print_help()
        sys.exit(1)

    # Read JSON spec
    if args.input == "-":
        raw = sys.stdin.read()
    else:
        input_path = Path(args.input)
        if not input_path.exists():
            print(f"ERROR: Input file not found: {args.input}", file=sys.stderr)
            sys.exit(1)
        raw = input_path.read_text(encoding="utf-8")

    try:
        spec = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON: {e}", file=sys.stderr)
        sys.exit(1)

    # Override output file if specified on CLI
    if args.output:
        spec["output_file"] = args.output

    # Determine base directory for image resolution
    base_dir = None
    if args.base_dir:
        base_dir = Path(args.base_dir).resolve()
    elif args.input and args.input != "-":
        base_dir = Path(args.input).resolve().parent

    title = spec.get("title", "Untitled")
    print(f"Building presentation: {title}")
    output_file = build_presentation(spec, base_dir=base_dir)
    print(f"Presentation saved to: {output_file}")


if __name__ == "__main__":
    main()
