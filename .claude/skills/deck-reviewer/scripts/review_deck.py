#!/usr/bin/env python3
"""
Deck Reviewer — automated visual-quality analysis for PPTX files.

Checks:
  1. Shape overlap detection
  2. Alignment & centering issues
  3. WCAG 2.1 color-contrast compliance (AA / AAA)
  4. Slide-type variety
  5. Line and connector clarity
  6. Visual density / clutter

Usage:
    python review_deck.py <path-to.pptx> [--json] [--out <report.json>]

Outputs a structured report (human-readable by default, JSON with --json).
"""

import argparse
import json
import math
import sys
from collections import Counter, defaultdict
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
except ImportError:
    sys.exit("python-pptx is required: pip install python-pptx")


# ── Helpers ──────────────────────────────────────────────────────────────────

def emu_to_inches(emu):
    """Convert EMUs to inches."""
    if emu is None:
        return 0.0
    return emu / 914400


def rgb_to_tuple(rgb_color):
    """Extract (R, G, B) 0-255 from an RGBColor or hex string."""
    if isinstance(rgb_color, RGBColor):
        return (rgb_color[0], rgb_color[1], rgb_color[2])
    if isinstance(rgb_color, str):
        c = rgb_color.lstrip("#")
        return tuple(int(c[i : i + 2], 16) for i in (0, 2, 4))
    return None


def srgb_to_linear(c):
    """Convert sRGB channel (0-255) to linear."""
    c = c / 255.0
    return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4


def relative_luminance(r, g, b):
    """WCAG 2.1 relative luminance."""
    return 0.2126 * srgb_to_linear(r) + 0.7152 * srgb_to_linear(g) + 0.0722 * srgb_to_linear(b)


def contrast_ratio(rgb1, rgb2):
    """WCAG contrast ratio between two (R,G,B) tuples."""
    l1 = relative_luminance(*rgb1)
    l2 = relative_luminance(*rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def box_contains(outer, inner):
    """True if outer fully contains inner (with small tolerance)."""
    tol = 0.1  # inches
    return (outer[0] - tol <= inner[0] and outer[1] - tol <= inner[1] and
            outer[2] + tol >= inner[2] and outer[3] + tol >= inner[3])


def boxes_overlap_partial(a, b):
    """True if two boxes partially overlap but neither contains the other.
    Returns False for full containment (intentional nesting) and trivial overlaps."""
    x_overlap = max(0, min(a[2], b[2]) - max(a[0], b[0]))
    y_overlap = max(0, min(a[3], b[3]) - max(a[1], b[1]))
    overlap_area = x_overlap * y_overlap
    if overlap_area == 0:
        return False
    area_a = (a[2] - a[0]) * (a[3] - a[1])
    area_b = (b[2] - b[0]) * (b[3] - b[1])
    if area_a == 0 or area_b == 0:
        return False
    # If one shape contains the other, this is intentional nesting (card pattern)
    if box_contains(a, b) or box_contains(b, a):
        return False
    # Only flag if overlap is significant relative to the smaller shape
    min_area = min(area_a, area_b)
    return (overlap_area / min_area) > 0.15


def bounding_box(shape):
    """Return (left, top, right, bottom) in inches, or None if missing geometry."""
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    l = emu_to_inches(shape.left)
    t = emu_to_inches(shape.top)
    return (l, t, l + emu_to_inches(shape.width), t + emu_to_inches(shape.height))


def classify_slide(shapes, layout_name):
    """Heuristic slide-type classification."""
    text_count = 0
    image_count = 0
    chart_count = 0
    table_count = 0
    shape_count = 0
    group_count = 0
    total_text_len = 0

    for sp in shapes:
        try:
            st = sp.shape_type
        except Exception:
            st = None

        if st == MSO_SHAPE_TYPE.PICTURE or st == MSO_SHAPE_TYPE.LINKED_PICTURE:
            image_count += 1
        elif st == MSO_SHAPE_TYPE.CHART or st == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            chart_count += 1
        elif st == MSO_SHAPE_TYPE.TABLE:
            table_count += 1
        elif st == MSO_SHAPE_TYPE.GROUP:
            group_count += 1
        else:
            shape_count += 1

        if sp.has_text_frame:
            text_count += 1
            total_text_len += len(sp.text_frame.text.strip())

    ln = (layout_name or "").lower()

    if "title" in ln and text_count <= 3 and image_count == 0:
        return "title"
    if image_count >= 2:
        return "image-gallery"
    if image_count == 1 and text_count <= 2:
        return "image-focus"
    if chart_count > 0:
        return "chart"
    if table_count > 0:
        return "table"
    if group_count >= 2 or shape_count >= 6:
        return "diagram"
    if text_count >= 3 and total_text_len > 300:
        return "text-heavy"
    if text_count == 1 and total_text_len < 80:
        return "statement"
    if total_text_len == 0 and shape_count == 0 and image_count == 0:
        return "blank"
    return "mixed-content"


def extract_text_colors(shape):
    """Yield (font_size_pt, fg_rgb_tuple, shape_name) for each run in shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            font = run.font
            size_pt = font.size.pt if font.size else None
            fg = None
            try:
                if font.color and font.color.rgb:
                    fg = rgb_to_tuple(font.color.rgb)
            except (AttributeError, TypeError):
                pass
            if fg:
                yield (size_pt, fg, shape.name if hasattr(shape, "name") else "unknown")


def extract_shape_fill(shape):
    """Try to get the solid-fill color of a shape. Returns (R,G,B) or None."""
    try:
        fill = shape.fill
        if fill.type is not None:
            from pptx.enum.dml import MSO_THEME_COLOR
            try:
                if fill.fore_color and fill.fore_color.rgb:
                    return rgb_to_tuple(fill.fore_color.rgb)
            except (AttributeError, TypeError):
                pass
    except Exception:
        pass
    return None


def extract_line_info(shape):
    """Return dict with line properties, or None."""
    try:
        ln = shape.line
        if ln is None:
            return None
        width_pt = ln.width.pt if ln.width else None
        color = None
        try:
            if ln.color and ln.color.rgb:
                color = rgb_to_tuple(ln.color.rgb)
        except (AttributeError, TypeError):
            pass
        if width_pt is not None or color is not None:
            return {"width_pt": width_pt, "color": color, "shape_name": shape.name}
    except Exception:
        pass
    return None


# ── Analysis passes ──────────────────────────────────────────────────────────

def check_overlaps(shapes):
    """Detect overlapping shapes, filtering out intentional patterns. Returns list of issue dicts."""
    issues = []
    bboxes = []
    for sp in shapes:
        bb = bounding_box(sp)
        if bb:
            bboxes.append((sp.name if hasattr(sp, "name") else "shape", bb, sp))

    for i in range(len(bboxes)):
        for j in range(i + 1, len(bboxes)):
            name_a, bb_a, sp_a = bboxes[i]
            name_b, bb_b, sp_b = bboxes[j]
            if not boxes_overlap_partial(bb_a, bb_b):
                continue

            # --- Filter out intentional overlap patterns ---

            # Text-on-image (standard design)
            a_is_text = sp_a.has_text_frame and sp_a.text_frame.text.strip()
            b_is_text = sp_b.has_text_frame and sp_b.text_frame.text.strip()
            try:
                a_is_image = sp_a.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)
            except Exception:
                a_is_image = False
            try:
                b_is_image = sp_b.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)
            except Exception:
                b_is_image = False
            if (a_is_text and b_is_image) or (b_is_text and a_is_image):
                continue

            # Chart + label (chart labels always overlap the chart area)
            try:
                a_is_chart = sp_a.shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT)
            except Exception:
                a_is_chart = False
            try:
                b_is_chart = sp_b.shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT)
            except Exception:
                b_is_chart = False
            if (a_is_chart and b_is_text) or (b_is_chart and a_is_text):
                continue

            # Adjacent text boxes that barely touch (common in multi-line layouts)
            if a_is_text and b_is_text:
                x_overlap = max(0, min(bb_a[2], bb_b[2]) - max(bb_a[0], bb_b[0]))
                y_overlap = max(0, min(bb_a[3], bb_b[3]) - max(bb_a[1], bb_b[1]))
                overlap_area = x_overlap * y_overlap
                smaller = min((bb_a[2]-bb_a[0])*(bb_a[3]-bb_a[1]),
                              (bb_b[2]-bb_b[0])*(bb_b[3]-bb_b[1]))
                if smaller > 0 and overlap_area / smaller < 0.25:
                    continue  # Minor text overlap — usually just generous bounding boxes

            issues.append({
                "type": "overlap",
                "severity": "warning",
                "shapes": [name_a, name_b],
                "message": f"Shapes '{name_a}' and '{name_b}' partially overlap — may cause visual confusion.",
            })
    return issues


def check_alignment(shapes, slide_width_in, slide_height_in):
    """Detect alignment/centering issues. Returns list of issue dicts."""
    issues = []
    NEAR_THRESHOLD = 0.15  # inches — close enough to look accidental
    center_x = slide_width_in / 2
    center_y = slide_height_in / 2

    bboxes = []
    for sp in shapes:
        bb = bounding_box(sp)
        if bb:
            bboxes.append((sp.name if hasattr(sp, "name") else "shape", bb))

    # Check near-center misses: shapes that are close to centered but not quite
    for name, bb in bboxes:
        shape_cx = (bb[0] + bb[2]) / 2
        shape_cy = (bb[1] + bb[3]) / 2
        dx = abs(shape_cx - center_x)
        dy = abs(shape_cy - center_y)
        shape_w = bb[2] - bb[0]

        # Only flag shapes that span >40% of slide width (likely meant to be centered)
        if shape_w > slide_width_in * 0.4:
            if 0.02 < dx < NEAR_THRESHOLD:
                issues.append({
                    "type": "alignment",
                    "severity": "info",
                    "shape": name,
                    "message": f"'{name}' is nearly centered horizontally (off by {dx:.2f}\") — consider centering exactly.",
                })

    # Check groups of shapes that are nearly aligned with each other
    lefts = [(name, bb[0]) for name, bb in bboxes]
    rights = [(name, bb[2]) for name, bb in bboxes]
    tops = [(name, bb[1]) for name, bb in bboxes]

    for edge_list, direction in [(lefts, "left"), (rights, "right"), (tops, "top")]:
        sorted_edges = sorted(edge_list, key=lambda x: x[1])
        for i in range(len(sorted_edges) - 1):
            n1, v1 = sorted_edges[i]
            n2, v2 = sorted_edges[i + 1]
            diff = abs(v2 - v1)
            if 0.02 < diff < NEAR_THRESHOLD:
                issues.append({
                    "type": "alignment",
                    "severity": "info",
                    "shapes": [n1, n2],
                    "message": f"'{n1}' and '{n2}' have nearly identical {direction} edges (off by {diff:.2f}\") — align them.",
                })
    return issues


def find_backing_fill(text_shape, shapes):
    """Find the fill color of a shape behind this text shape (scrim/rectangle pattern).

    Walks backwards through the shape list (z-order) looking for a filled shape
    whose bounding box covers the text shape. This detects scrim overlays and
    colored rectangles behind text boxes.
    """
    text_bb = bounding_box(text_shape)
    if not text_bb:
        return None

    # Build list with indices to check shapes that appear before (behind) this one
    shape_list = list(shapes)
    try:
        text_idx = shape_list.index(text_shape)
    except ValueError:
        return None

    # Walk backwards from the text shape to find a backing fill
    for i in range(text_idx - 1, -1, -1):
        candidate = shape_list[i]
        cand_bb = bounding_box(candidate)
        if not cand_bb:
            continue
        # Check if candidate covers the text shape
        if box_contains(cand_bb, text_bb):
            fill = extract_shape_fill(candidate)
            if fill:
                return fill
    return None


def check_contrast(shapes, slide_bg_rgb):
    """Check WCAG contrast compliance. Returns list of issue dicts."""
    issues = []
    bg = slide_bg_rgb or (255, 255, 255)  # default white
    shape_list = list(shapes)

    for sp in shape_list:
        # Determine effective background: shape's own fill > backing shape fill > slide bg
        shape_fill = extract_shape_fill(sp)
        if shape_fill:
            effective_bg = shape_fill
        else:
            backing = find_backing_fill(sp, shape_list)
            effective_bg = backing if backing else bg

        for size_pt, fg_rgb, sp_name in extract_text_colors(sp):
            ratio = contrast_ratio(fg_rgb, effective_bg)
            is_large = (size_pt and size_pt >= 18) or (size_pt and size_pt >= 14)  # simplified
            threshold_aa = 3.0 if is_large else 4.5
            threshold_aaa = 4.5 if is_large else 7.0

            if ratio < threshold_aa:
                issues.append({
                    "type": "contrast",
                    "severity": "error",
                    "shape": sp_name,
                    "fg_color": f"#{fg_rgb[0]:02x}{fg_rgb[1]:02x}{fg_rgb[2]:02x}",
                    "bg_color": f"#{effective_bg[0]:02x}{effective_bg[1]:02x}{effective_bg[2]:02x}",
                    "ratio": round(ratio, 2),
                    "required_aa": threshold_aa,
                    "message": f"WCAG AA FAIL in '{sp_name}': contrast {ratio:.2f}:1 (needs {threshold_aa}:1). "
                               f"Text #{fg_rgb[0]:02x}{fg_rgb[1]:02x}{fg_rgb[2]:02x} on #{effective_bg[0]:02x}{effective_bg[1]:02x}{effective_bg[2]:02x}.",
                })
            elif ratio < threshold_aaa:
                issues.append({
                    "type": "contrast",
                    "severity": "warning",
                    "shape": sp_name,
                    "fg_color": f"#{fg_rgb[0]:02x}{fg_rgb[1]:02x}{fg_rgb[2]:02x}",
                    "bg_color": f"#{effective_bg[0]:02x}{effective_bg[1]:02x}{effective_bg[2]:02x}",
                    "ratio": round(ratio, 2),
                    "required_aaa": threshold_aaa,
                    "message": f"WCAG AAA FAIL in '{sp_name}': contrast {ratio:.2f}:1 (needs {threshold_aaa}:1). "
                               f"Passes AA but not AAA.",
                })
    return issues


def check_lines(shapes):
    """Check line clarity. Returns list of issue dicts."""
    issues = []
    for sp in shapes:
        info = extract_line_info(sp)
        if info and info["width_pt"] is not None:
            if info["width_pt"] < 0.5:
                issues.append({
                    "type": "line-clarity",
                    "severity": "warning",
                    "shape": info["shape_name"],
                    "width_pt": info["width_pt"],
                    "message": f"Line on '{info['shape_name']}' is very thin ({info['width_pt']:.2f}pt) — may be invisible when projected.",
                })
            if info["color"]:
                # Light gray lines on white are hard to see
                bg = (255, 255, 255)
                ratio = contrast_ratio(info["color"], bg)
                if ratio < 2.0:
                    issues.append({
                        "type": "line-clarity",
                        "severity": "warning",
                        "shape": info["shape_name"],
                        "message": f"Line on '{info['shape_name']}' has very low contrast ({ratio:.2f}:1 vs white) — may be invisible.",
                    })
    return issues


def check_density(shapes, slide_width_in, slide_height_in):
    """Check visual density / clutter. Returns list of issue dicts."""
    issues = []
    slide_area = slide_width_in * slide_height_in
    shape_count = 0
    covered_area = 0.0

    for sp in shapes:
        bb = bounding_box(sp)
        if bb:
            shape_count += 1
            w = bb[2] - bb[0]
            h = bb[3] - bb[1]
            covered_area += w * h

    if shape_count > 15:
        issues.append({
            "type": "density",
            "severity": "warning",
            "shape_count": shape_count,
            "message": f"Slide has {shape_count} shapes — consider simplifying for readability.",
        })

    coverage = covered_area / slide_area if slide_area > 0 else 0
    if coverage > 0.9:
        issues.append({
            "type": "density",
            "severity": "info",
            "coverage_pct": round(coverage * 100, 1),
            "message": f"Shape bounding boxes cover {coverage * 100:.0f}% of slide area — may feel cramped.",
        })

    return issues


def check_variety(slide_types):
    """Check slide-type variety across the deck. Returns list of issue dicts."""
    issues = []
    total = len(slide_types)
    if total == 0:
        return issues

    counts = Counter(slide_types)
    unique_types = len(counts)

    if total >= 5 and unique_types < 3:
        issues.append({
            "type": "variety",
            "severity": "warning",
            "unique_types": unique_types,
            "total_slides": total,
            "distribution": dict(counts),
            "message": f"Only {unique_types} distinct slide type(s) across {total} slides — add variety (images, charts, diagrams, statements) to keep the audience engaged.",
        })

    # Flag dominant type
    most_common_type, most_common_count = counts.most_common(1)[0]
    dominance = most_common_count / total
    if total >= 5 and dominance > 0.6:
        issues.append({
            "type": "variety",
            "severity": "info",
            "dominant_type": most_common_type,
            "dominance_pct": round(dominance * 100, 1),
            "message": f"'{most_common_type}' slides make up {dominance * 100:.0f}% of the deck — consider breaking up the rhythm.",
        })

    return issues


# ── Slide background extraction ──────────────────────────────────────────────

def get_slide_bg(slide):
    """Try to determine the slide background color. Returns (R,G,B) or None."""
    try:
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            try:
                if bg.fill.fore_color and bg.fill.fore_color.rgb:
                    return rgb_to_tuple(bg.fill.fore_color.rgb)
            except Exception:
                pass
    except Exception:
        pass
    return None


# ── Main analysis ────────────────────────────────────────────────────────────

def analyze_deck(pptx_path):
    """Run all checks on a PPTX file. Returns structured report dict."""
    prs = Presentation(pptx_path)
    slide_width_in = emu_to_inches(prs.slide_width)
    slide_height_in = emu_to_inches(prs.slide_height)

    report = {
        "file": str(pptx_path),
        "slide_count": len(prs.slides),
        "slide_dimensions": f"{slide_width_in:.2f}\" x {slide_height_in:.2f}\"",
        "slides": [],
        "summary": {
            "total_issues": 0,
            "errors": 0,
            "warnings": 0,
            "info": 0,
            "slide_types": {},
            "variety_issues": [],
            "contrast_failures": 0,
            "overlap_count": 0,
            "alignment_issues": 0,
        },
    }

    slide_types = []

    for idx, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else None
        shapes = list(slide.shapes)
        slide_bg = get_slide_bg(slide)
        slide_type = classify_slide(shapes, layout_name)
        slide_types.append(slide_type)

        all_issues = []
        all_issues.extend(check_overlaps(shapes))
        all_issues.extend(check_alignment(shapes, slide_width_in, slide_height_in))
        all_issues.extend(check_contrast(shapes, slide_bg))
        all_issues.extend(check_lines(shapes))
        all_issues.extend(check_density(shapes, slide_width_in, slide_height_in))

        slide_report = {
            "slide_number": idx,
            "layout": layout_name,
            "type": slide_type,
            "shape_count": len(shapes),
            "issues": all_issues,
        }
        report["slides"].append(slide_report)

        for issue in all_issues:
            sev = issue.get("severity", "info")
            report["summary"]["total_issues"] += 1
            if sev == "error":
                report["summary"]["errors"] += 1
            elif sev == "warning":
                report["summary"]["warnings"] += 1
            else:
                report["summary"]["info"] += 1
            if issue["type"] == "contrast" and sev == "error":
                report["summary"]["contrast_failures"] += 1
            if issue["type"] == "overlap":
                report["summary"]["overlap_count"] += 1
            if issue["type"] == "alignment":
                report["summary"]["alignment_issues"] += 1

    # Deck-wide variety check
    variety_issues = check_variety(slide_types)
    report["summary"]["variety_issues"] = variety_issues
    report["summary"]["slide_types"] = dict(Counter(slide_types))
    for issue in variety_issues:
        report["summary"]["total_issues"] += 1
        sev = issue.get("severity", "info")
        if sev == "warning":
            report["summary"]["warnings"] += 1
        else:
            report["summary"]["info"] += 1

    return report


def format_report(report):
    """Return a human-readable string from the report dict."""
    lines = []
    lines.append(f"=== Deck Review: {report['file']} ===")
    lines.append(f"Slides: {report['slide_count']}  |  Dimensions: {report['slide_dimensions']}")
    lines.append(f"Slide types: {report['summary']['slide_types']}")
    lines.append("")

    s = report["summary"]
    lines.append(f"--- Summary ---")
    lines.append(f"Total issues: {s['total_issues']}  (errors: {s['errors']}, warnings: {s['warnings']}, info: {s['info']})")
    lines.append(f"  WCAG contrast failures (AA): {s['contrast_failures']}")
    lines.append(f"  Overlapping shapes: {s['overlap_count']}")
    lines.append(f"  Alignment near-misses: {s['alignment_issues']}")
    lines.append("")

    if s["variety_issues"]:
        lines.append("--- Slide Variety ---")
        for vi in s["variety_issues"]:
            lines.append(f"  [{vi['severity'].upper()}] {vi['message']}")
        lines.append("")

    for sl in report["slides"]:
        if sl["issues"]:
            lines.append(f"--- Slide {sl['slide_number']} ({sl['type']}, layout: {sl['layout']}) ---")
            for issue in sl["issues"]:
                lines.append(f"  [{issue['severity'].upper()}] {issue['message']}")
            lines.append("")

    if s["total_issues"] == 0:
        lines.append("No issues detected. The deck looks clean!")

    return "\n".join(lines)


# ── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Review a PPTX deck for visual quality issues.")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    parser.add_argument("--out", help="Write report to file instead of stdout")
    args = parser.parse_args()

    path = Path(args.pptx)
    if not path.exists():
        sys.exit(f"File not found: {path}")
    if not path.suffix.lower() == ".pptx":
        sys.exit(f"Expected a .pptx file, got: {path.suffix}")

    report = analyze_deck(path)

    if args.json:
        output = json.dumps(report, indent=2)
    else:
        output = format_report(report)

    if args.out:
        Path(args.out).write_text(output)
        print(f"Report written to {args.out}")
    else:
        print(output)


if __name__ == "__main__":
    main()
